# Creating powerpoint presentations using the python-pptx package

from pptx import Presentation
from pptx.util import Inches

X = Presentation()

Layout = X.slide_layouts[0]
first_slide = X.slides.add_slide(Layout)

first_slide.shapes.title.text = "Creating a powerpoint using Python"
first_slide.placeholders[1].text = "CREATED BY ME?! :D"

X.save("First_presentation.pptx")


Second_Layout = X.slide_layouts[5]
second_slide = X.slides.add_slide(Second_Layout)
second_slide.shapes.title.text = "Second slide"

textbox = second_slide.shapes.add_textbox(Inches(3), Inches(1.5),Inches(3), Inches(1))
textframe = textbox.text_frame
paragraph = textframe.add_paragraph() 
paragraph.text = "This is a paragraph in the second slide!"

X.save("First_presentation.pptx")